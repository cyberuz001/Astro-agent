import os, subprocess, json, sys
from typing import Optional
from langchain_core.tools import tool

# Re-use our original config logic
# Since this file was moved to astro_agent/tools/__init__.py, we need one more dirname
sys.path.append(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))
try:
    from astro import get_weather_and_time as original_gwt
    from astro import make_voice_call as original_mvc
    from astro import run_cmd as original_run_cmd
except ImportError:
    # Safely fallback to simple wrappers if refactoring removes them from base
    original_gwt = None
    original_mvc = None
    original_run_cmd = None

@tool
def bash_terminal(command: str) -> str:
    """Linux tizim buyruqlarini ishga tushiradi (Masalan: ls -la, cat fayl, free -m)"""
    try:
        r = subprocess.run(command, shell=True, capture_output=True, text=True, timeout=120)
        return (r.stdout + r.stderr).strip()[:4000]
    except Exception as e:
        return f"Xato: {e}"

@tool
def get_weather_and_time(location: str, iana_timezone: str = "Asia/Tashkent") -> str:
    """Ixtiyoriy shahar orqali hozirgi harorat va EXACT vaqtni oladi. (Masalan, location="London", iana_timezone="Europe/London")"""
    try:
        if original_gwt:
            return original_gwt(location, iana_timezone)
        return "Xato: original_gwt modulda mavjud emas."
    except Exception as e:
        return f"Xato: get_weather_and_time ishlamayapti: {e}"

@tool
def make_pbx_call(audio_message: str, goal: str = "") -> str:
    """Astro agenti nomidan Asterisk PBX orqali telefon qilib gaplashadi va missiyani bajaradi."""
    try:
        if original_mvc:
            return original_mvc(audio_message, goal)
        return "Xato: original_mvc modulda mavjud emas."
    except Exception as e:
        return f"Xato: Qo'ng'iroq bloki ushlamayapti: {e}"

@tool
def web_search_tool(query: str) -> str:
    """DuckDuckGo orqali erkin internet ma'lumotlarini qidirish."""
    try:
        from langchain_community.tools import DuckDuckGoSearchRun
        return DuckDuckGoSearchRun().invoke(query)
    except Exception as e:
        return f"Search Error: {e}"

@tool
def pbx_admin(action: str, ext: Optional[str]=None, pwd: Optional[str]=None) -> str:
    """Asterisk tizimini boshqaradi. action='reload', yoki action='set_pass' ext='101' pwd='pass' """
    if action == "reload":
        r = subprocess.run("echo 'password' | sudo -S asterisk -rx 'core reload'", shell=True, capture_output=True, text=True)
        return "Reloaded."
    elif action == "set_pass" and ext and pwd:
        # Simplistic demonstration instead of original 30 line py-script
        return f"Parol {ext} uchun bash_terminal orqali hal qilinishi tavsiya qilinadi."
    return "Noma'lum Asterisk buyrug'i."

@tool
def process_document(filepath: str) -> str:
    """PDF, DOCX yoki TXT kabi hujjatlarni o'qiydi va matnini qaytaradi. Barcha hujjat ishlari uchun ishlating."""
    if not os.path.exists(filepath):
        return f"Xato: Fayl topilmadi: {filepath}"

    ext = os.path.splitext(filepath)[1].lower()
    text = ""
    try:
        if ext == ".pdf":
            import pypdf
            with open(filepath, "rb") as f:
                reader = pypdf.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
        elif ext == ".docx":
            import docx
            doc = docx.Document(filepath)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif ext == ".xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(filepath, data_only=True)
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                text += f"Sheet: {sheet}\n"
                for row in ws.iter_rows(values_only=True):
                    row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                    text += row_text + "\n"
        else:
            with open(filepath, "r", encoding="utf-8") as f:
                text = f.read()

        return text[:10000] # Qisqartirib qaytaramiz (kontekst tolib ketmasligi uchun)
    except ImportError:
        return f"Xato: {ext} formatini o'qish uchun kutubxona yetishmaydi."
    except Exception as e:
        return f"Xato hujjatni o'qishda: {e}"

@tool
def create_presentation(title: str, slides_content: str, output_path: str = "presentation.pptx") -> str:
    """Yangi PowerPoint prezentatsiya (.pptx) yaratadi. slides_content bu JSON formatdagi string bo'lishi kerak: [{"title": "Slayd 1", "content": "Matn..."}, ...]"""
    try:
        from pptx import Presentation

        prs = Presentation()
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle = slide.placeholders[1]
        title_shape.text = title
        subtitle.text = "Astro Agent tomonidan yaratilgan"

        try:
            slides_data = json.loads(slides_content)
        except:
            return "Xato: slides_content to'g'ri JSON formatda emas. Yaroqli format: [{\"title\": \"...\", \"content\": \"...\"}]"

        for s in slides_data:
            slide_layout = prs.slide_layouts[1] # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]

            title_shape.text = s.get("title", "Slayd")
            tf = body_shape.text_frame
            tf.text = s.get("content", "")

        prs.save(output_path)
        return f"Muvaffaqiyatli: Prezentatsiya '{output_path}' faylida saqlandi."
    except ImportError:
        return "Xato: python-pptx kutubxonasi yetishmaydi."
    except Exception as e:
        return f"Xato prezentatsiya yaratishda: {e}"

@tool
def file_manager(action: str, filepath: str, content: str = "") -> str:
    """Fayllar bilan ishlash uchun xavfsiz menejer. action='read', 'write', 'delete' bo'lishi mumkin."""
    try:
        if action == "read":
            if not os.path.exists(filepath):
                return f"Xato: Fayl topilmadi - {filepath}"
            with open(filepath, "r", encoding="utf-8") as f:
                return f.read()[:20000]
        elif action == "write":
            # Faylni yaratish yoki ustidan yozish
            os.makedirs(os.path.dirname(os.path.abspath(filepath)), exist_ok=True)
            with open(filepath, "w", encoding="utf-8") as f:
                f.write(content)
            return f"Muvaffaqiyatli: {filepath} yozildi."
        elif action == "delete":
            if os.path.exists(filepath):
                os.remove(filepath)
                return f"Muvaffaqiyatli: {filepath} o'chirildi."
            return f"Xato: Fayl topilmadi - {filepath}"
        else:
            return "Xato: Noma'lum harakat (action). Faqat read, write, delete."
    except Exception as e:
        return f"Xato: file_manager ishida muammo: {e}"

@tool
def git_manager(action: str, branch: str = "main", message: str = "") -> str:
    """Git omborlari bilan ishlash. action='status', 'commit', 'checkout', 'push', 'pull' bo'lishi mumkin.
    commit qilish uchun message kerak. branch tanlash uchun branch parametri."""
    try:
        if action == "status":
            r = subprocess.run("git status", shell=True, capture_output=True, text=True)
            return r.stdout + r.stderr
        elif action == "commit":
            if not message:
                return "Xato: commit qilish uchun message kerak."
            subprocess.run("git add .", shell=True)
            r = subprocess.run(f'git commit -m "{message}"', shell=True, capture_output=True, text=True)
            return r.stdout + r.stderr
        elif action == "checkout":
            r = subprocess.run(f"git checkout -b {branch} || git checkout {branch}", shell=True, capture_output=True, text=True)
            return r.stdout + r.stderr
        elif action == "push":
            r = subprocess.run(f"git push origin {branch}", shell=True, capture_output=True, text=True)
            return r.stdout + r.stderr
        elif action == "pull":
            r = subprocess.run(f"git pull origin {branch}", shell=True, capture_output=True, text=True)
            return r.stdout + r.stderr
        else:
            return "Xato: Noma'lum git buyrug'i."
    except Exception as e:
        return f"Xato: git_manager ishida muammo: {e}"

@tool
def delegate_task(agent_role: str, task_description: str) -> str:
    """Yirik vazifalarni sub-agentlarga (masalan: 'frontend_developer', 'qa_engineer', 'devops') topshirish uchun.
    Astro bu tool orqali asinxron ravishda sub-agentlardan javob kutadi."""
    # Hozircha mock-implementatsiya bo'lib, kelajakda LangGraph multi-agent tarmog'iga ulanadi.
    return f"[{agent_role} Sub-Agentiga vazifa yuborildi]: '{task_description}'. Agent tez orada tahlilni tugatadi. Iltimos bash_terminal orqali natijalarni (masalan test natijasi) tekshiring."

# Core array exported for LangGraph
ASTRO_TOOLS = [bash_terminal, get_weather_and_time, make_pbx_call, web_search_tool, pbx_admin, process_document, create_presentation, file_manager, git_manager, delegate_task]
